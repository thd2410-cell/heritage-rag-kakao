# -*- coding: utf-8 -*-
"""부록 진입 다리 페이지: QR(웹 체험) + 부록 안내 + 감사 클로징. 앞 양식(Noto Sans KR/네이비)."""
import qrcode
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SRC = r"c:\Users\SSAFY\Desktop\heritage-rag-kakao\data_eng\presentation\유산톡 AI (2)_간트_v2.pptx"
DST = r"c:\Users\SSAFY\Desktop\heritage-rag-kakao\data_eng\presentation\유산톡 AI (2)_간트_v3.pptx"
QR_PATH = r"c:\Users\SSAFY\Desktop\heritage-rag-kakao\data_eng\presentation\_qr_web.png"

FONT="Noto Sans KR"
NAVY=RGBColor(0x1A,0x2C,0x5B); GRAY=RGBColor(0x71,0x80,0x96); BODY=RGBColor(0x4A,0x55,0x68)
WHITE=RGBColor(0xFF,0xFF,0xFF); YELLOW=RGBColor(0xF2,0xC8,0x00)
LIGHTNAVY=RGBColor(0xC9,0xD3,0xE6)

# --- QR 생성 (네이비) ---
qr = qrcode.QRCode(box_size=14, border=1)
qr.add_data("https://heritage-chat.com"); qr.make(fit=True)
qr.make_image(fill_color="#1A2C5B", back_color="white").save(QR_PATH)

prs = Presentation(SRC)
blank = prs.slide_layouts[0]

def add_slide():
    s = prs.slides.add_slide(blank)
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    return s

def set_font(run, size, color, bold=False):
    run.font.name=FONT; run.font.size=Pt(size); run.font.color.rgb=color; run.font.bold=bold
    rPr=run._r.get_or_add_rPr(); ea=rPr.find(qn('a:ea'))
    if ea is None: ea=rPr.makeelement(qn('a:ea'),{}); rPr.append(ea)
    ea.set('typeface', FONT)

def textbox(slide,l,t,w,h,text,size,color,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    tb=slide.shapes.add_textbox(Emu(l),Emu(t),Emu(w),Emu(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    tf.vertical_anchor=anchor
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=text; set_font(r,size,color,bold)
    return tb

def rrect(slide,l,t,w,h,fill,rounded=True):
    shp=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                               Emu(l),Emu(t),Emu(w),Emu(h))
    shp.fill.solid(); shp.fill.fore_color.rgb=fill; shp.line.fill.background()
    shp.shadow.inherit=False
    return shp

def header(slide,title,subtitle):
    bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Emu(571500),Emu(420000),Emu(45720),Emu(300000))
    bar.fill.solid(); bar.fill.fore_color.rgb=NAVY; bar.line.fill.background(); bar.shadow.inherit=False
    textbox(slide,723900,381000,10772775,430000,title,22.5,NAVY,bold=True)
    textbox(slide,723900,885825,10972800,240000,subtitle,11.5,GRAY)
    div=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Emu(571500),Emu(1230000),Emu(11049000),Emu(22860))
    div.fill.solid(); div.fill.fore_color.rgb=YELLOW; div.line.fill.background(); div.shadow.inherit=False

# ================= 다리 페이지 =================
s = add_slide()
header(s, "직접 확인해보세요",
       "유산잇다는 '기획' 속 약속이 아니라, 지금 동작하는 서비스입니다")

# --- 왼쪽: QR 체험 ---
# QR 카드 배경
rrect(s, 1180000, 1700000, 4350000, 3650000, RGBColor(0xF4,0xF6,0xFA))
qr_sz = 2550000
s.shapes.add_picture(QR_PATH, Emu(int(1180000+(4350000-qr_sz)/2)), Emu(2050000), Emu(qr_sz), Emu(qr_sz))
textbox(s, 1180000, 4720000, 4350000, 280000, "웹에서 바로 체험", 15, NAVY, bold=True, align=PP_ALIGN.CENTER)
textbox(s, 1180000, 5010000, 4350000, 240000, "heritage-chat.com", 12, BODY, align=PP_ALIGN.CENTER)
textbox(s, 1180000, 5260000, 4350000, 220000, "카카오톡 채널에서도 지금 대화할 수 있습니다", 9.5, GRAY, align=PP_ALIGN.CENTER)
# "지금 스캔" 뱃지
badge = rrect(s, 4520000, 1880000, 1020000, 330000, YELLOW)
textbox(s, 4520000, 1880000, 1020000, 330000, "지금 스캔 ▶", 10, NAVY, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# --- 오른쪽: 부록 안내 패널 ---
panel = rrect(s, 6050000, 1700000, 5570000, 3650000, NAVY)
tb = s.shapes.add_textbox(Emu(6400000), Emu(1960000), Emu(4900000), Emu(3200000))
tf = tb.text_frame; tf.word_wrap=True
tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
def para(text, size, color, bold=False, space_before=8):
    p = tf.add_paragraph() if tf.paragraphs[0].runs else tf.paragraphs[0]
    p.space_before = Pt(space_before)
    r=p.add_run(); r.text=text; set_font(r,size,color,bold); return p
para("더 깊은 근거는 부록에서", 15, WHITE, bold=True, space_before=0)
para("발표에서 다 담지 못한 상세 설계·검증을 정리했습니다.", 10, LIGHTNAVY, space_before=4)
for t in ["부록 A — 상세 기능 명세 (15)",
          "부록 B — 상세 ERD (15개 테이블)",
          "부록 C — KPI 체계",
          "부록 D — KPT 개발 인사이트",
          "부록 E — 예상 Q&A",
          "개발 현황 · 추진 간트 · 30분 상세 계획"]:
    para("·  "+t, 11.5, WHITE, bold=False, space_before=9)

# --- 클로징 바 ---
rrect(s, 571500, 5650000, 11049000, 720000, YELLOW)
ctb = s.shapes.add_textbox(Emu(571500), Emu(5650000), Emu(11049000), Emu(720000))
ctf = ctb.text_frame; ctf.word_wrap=True; ctf.vertical_anchor=MSO_ANCHOR.MIDDLE
ctf.margin_left=Emu(200000); ctf.margin_right=Emu(200000)
cp=ctf.paragraphs[0]; cp.alignment=PP_ALIGN.CENTER
r=cp.add_run(); r.text="기획서 속 약속이 아니라, 지금 카카오톡에서 답하는 서비스입니다.  끝까지 살펴봐 주셔서 감사합니다."
set_font(r, 13.5, NAVY, bold=True)

prs.save(DST)
print("saved:", DST, "| total:", len(prs.slides._sldIdLst))
