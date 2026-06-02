# -*- coding: utf-8 -*-
"""유산톡 AI (2).pptx 뒤에 간트(현황+계획) 2장 추가. 앞 양식(Noto Sans KR/네이비) 맞춤."""
import copy
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SRC = r"c:\Users\SSAFY\Desktop\heritage-rag-kakao\data_eng\presentation\유산톡 AI (2).pptx"
DST = r"c:\Users\SSAFY\Desktop\heritage-rag-kakao\data_eng\presentation\유산톡 AI (2)_간트_v2.pptx"

FONT = "Noto Sans KR"
NAVY = RGBColor(0x1A, 0x2C, 0x5B)
GRAY = RGBColor(0x71, 0x80, 0x96)
BODY = RGBColor(0x4A, 0x55, 0x68)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW = RGBColor(0xF2, 0xC8, 0x00)
# status tints
G0 = (RGBColor(0xD6,0xF5,0xDD), RGBColor(0x1B,0x5E,0x20))   # 완료(green)
G1 = (RGBColor(0xDC,0xEB,0xFB), RGBColor(0x12,0x49,0x8A))   # 사전(blue)
G2 = (RGBColor(0xFF,0xF3,0xCC), RGBColor(0x7A,0x5C,0x00))   # 현장(amber)

prs = Presentation(SRC)
blank = prs.slide_layouts[0]

def add_slide():
    s = prs.slides.add_slide(blank)
    # remove inherited placeholders
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    # white background
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    return s

def set_font(run, size, color, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    # set east-asian font too
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', FONT)

def textbox(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    set_font(r, size, color, bold)
    return tb

def header(slide, title, subtitle):
    # accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(571500), Emu(420000), Emu(45720), Emu(300000))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    textbox(slide, 723900, 381000, 10772775, 430000, title, 22.5, NAVY, bold=True)
    textbox(slide, 723900, 885825, 10972800, 240000, subtitle, 11.5, GRAY)
    # yellow divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(571500), Emu(1230000), Emu(11049000), Emu(22860))
    div.fill.solid(); div.fill.fore_color.rgb = YELLOW; div.line.fill.background()

def style_cell(cell, text, size, color, fill, bold=False, align=PP_ALIGN.LEFT):
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.margin_left = Emu(45720); cell.margin_right = Emu(45720)
    cell.margin_top = Emu(18000); cell.margin_bottom = Emu(18000)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame; tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    set_font(r, size, color, bold)

def make_table(slide, rows, cols, l, t, w, h, col_w):
    gf = slide.shapes.add_table(rows, cols, Emu(l), Emu(t), Emu(w), Emu(h))
    tbl = gf.table
    # disable banded/first-row style emphasis by clearing style id -> keep simple
    tbl.first_row = False; tbl.horz_banding = False
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = Emu(cw)
    return tbl

# ============== SLIDE 1: 현황 + 추진 일정 (간트 매트릭스) ==============
s1 = add_slide()
header(s1,
       "현재 구현 현황을 기반으로 본선까지 단계별로 완성한다",
       "예선 프로토타입(완료) → 사전 개발(Sprint 1) → 본선 현장(Sprint 2)의 3단계 추진 로드맵")

# milestones line
textbox(s1, 571500, 1320000, 11049000, 230000,
        "◆ 5/31 예선 기획서 제출(현재)      ◆ 6/5 예선 결과      ◆ 6/13 본선 시작      ◆ 6/14 발표",
        11, NAVY, bold=True)

cols = ["워크스트림", "Phase 0 · 완료 (~5/31)", "Sprint 1 · 사전개발 (6/5~6/12)", "Sprint 2 · 본선현장 (6/13~14)"]
rows_data = [
 ("데이터·임베딩", "17,840건 수집·정제·bge-m3 임베딩·facet_json 생성", "좌표·행사 API 연동, 전체 재임베딩, alias 사전 확장", "모르는 유산 fallback 보완"),
 ("AI·RAG", "하이브리드 라우터·가드레일·80문항 검증(단일 Recall≈100%·충실성 0.89)", "6브랜치 라우터 완성, LLM Soft Filter, 프롬프트 보완", "비교질문(compare), 응답 2차 튜닝"),
 ("백엔드·카카오", "FastAPI·Kakao Skill·pgvector — 프로토타입 V1.0 체험 가능", "말풍선 제어, 멀티턴 강화, BasicCard/Carousel", "컨텍스트 유지, 초기화 커맨드"),
 ("게이미피케이션", "퀴즈·카드 로직 1차 구현(고도화중)", "동적 퀴즈 생성, 카드·도감 WebView 완성", "피드백 수집(👍👎), 오늘의 유산"),
 ("인프라·QA", "Docker·Cloudflare Tunnel 배포 완료", "통합 테스트, 배포 최종 확정", "서버 안정화·부하 테스트"),
 ("발표·기획", "예선 기획서 제출 완료", "—", "발표자료·데모 시나리오·리허설"),
]
nrows = len(rows_data) + 1
top = 1620000
tbl = make_table(s1, nrows, 4, 571500, top, 11049000, 4400000,
                 [1900000, 3300000, 3149000, 2700000])
# header row
for j, c in enumerate(cols):
    style_cell(tbl.cell(0, j), c, 10, WHITE, NAVY, bold=True,
               align=PP_ALIGN.CENTER if j else PP_ALIGN.CENTER)
phase_fill = [None, G0, G1, G2]
for i, rd in enumerate(rows_data, start=1):
    style_cell(tbl.cell(i, 0), rd[0], 9.5, WHITE, NAVY, bold=True, align=PP_ALIGN.CENTER)
    for j in range(1, 4):
        fill, txtcol = phase_fill[j]
        style_cell(tbl.cell(i, j), rd[j], 8.5, txtcol, fill)

# legend
textbox(s1, 571500, 6120000, 11049000, 230000,
        "■ 완료(Phase 0)   ■ 사전개발(Sprint 1)   ■ 본선현장(Sprint 2)      ※ 현재 작동하는 프로토타입(heritage-chat.com) 위에 기능을 단계적으로 확장",
        9, GRAY)

# ============== SLIDE 2: 진짜 간트차트 (시간축 막대) ==============
def add_rect(slide, l, t, w, h, fill, line_color=None, line_w=0.75):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(l)), Emu(int(t)), Emu(int(w)), Emu(int(h)))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line_color is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_color; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

sg = add_slide()
header(sg,
       "예선 제출부터 본선 발표까지, 한눈에 보는 추진 간트",
       "완료 영역(현재) 위에 Sprint 1·2 작업을 일정에 배치한 전체 개발 로드맵")

CHART_L = 2650000; CHART_T = 2150000; ROW_H = 545000; NROW = 6
CHART_R = 11620500; CHART_W = CHART_R - CHART_L; SEG = 11; SEG_W = CHART_W/SEG
CHART_B = CHART_T + ROW_H*NROW
def bx(i): return CHART_L + i*SEG_W
GRID = RGBColor(0xE2,0xE8,0xF0)
DONE = (RGBColor(0xC8,0xEA,0xD4), RGBColor(0x1B,0x5E,0x20))
S1c = (RGBColor(0xCF,0xE2,0xF7), RGBColor(0x12,0x49,0x8A))
S2c = (RGBColor(0xFC,0xE9,0xB8), RGBColor(0x7A,0x5C,0x00))
STAT = {"done":DONE, "s1":S1c, "s2":S2c}

# horizontal row separators + vertical day gridlines
for ri in range(NROW+1):
    add_rect(sg, CHART_L, CHART_T+ri*ROW_H-5715, CHART_W, 11430, GRID)
for i in range(SEG+1):
    add_rect(sg, bx(i)-5715, CHART_T, 11430, ROW_H*NROW, GRID)
# date labels
date_marks = [(0,"완료(~현재)",False,DONE[1]),(1,"6/5",True,NAVY),(3,"6/7",False,GRAY),
              (5,"6/9",False,GRAY),(7,"6/11",False,GRAY),(9,"6/13",True,NAVY),(10,"6/14",True,NAVY)]
for i,txt,ms,col in date_marks:
    cx = bx(i) if i!=0 else bx(0)+SEG_W/2
    tb = textbox(sg, int(cx)-450000, 1830000, 900000, 200000, txt, 8.5, col, bold=ms, align=PP_ALIGN.CENTER)
# milestone vertical lines + ◆ labels
for i,txt in [(1,"◆예선결과"),(9,"◆본선"),(10,"◆발표")]:
    add_rect(sg, bx(i)-9000, CHART_T, 18000, ROW_H*NROW, NAVY)
    textbox(sg, int(bx(i))-350000, 1630000, 700000, 180000, txt, 7.5, NAVY, bold=True, align=PP_ALIGN.CENTER)
# rows + bars
gantt_rows = [
 ("데이터·임베딩", [(0,1,"done","수집·정제·bge-m3·facet"),(2,4,"s1","재임베딩·좌표/행사"),(9,10,"s2","fallback")]),
 ("AI·RAG", [(0,1,"done","라우터·검증 0.89"),(4,6,"s1","6브랜치·SoftFilter"),(9,11,"s2","비교·튜닝")]),
 ("백엔드·카카오", [(0,1,"done","FastAPI·Skill·V1.0"),(6,8,"s1","말풍선·멀티턴·카드"),(9,10,"s2","컨텍스트·초기화")]),
 ("게이미피케이션", [(0,1,"done","퀴즈/카드 로직"),(6,8,"s1","동적퀴즈·도감"),(9,11,"s2","피드백·오늘유산")]),
 ("인프라·QA", [(0,1,"done","Docker·Cloudflare"),(8,9,"s1","통합·배포"),(10,11,"s2","안정화·부하")]),
 ("발표·기획", [(9,11,"s2","발표자료·리허설·발표")]),
]
BAR_H = int(ROW_H*0.58)
for ri,(label, bars) in enumerate(gantt_rows):
    rt = CHART_T + ri*ROW_H
    rl = textbox(sg, 571500, rt, 2030000, ROW_H, label, 9, NAVY, bold=True)
    rl.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for (si,ei,st,txt) in bars:
        fill,tcol = STAT[st]
        x0 = bx(si); w = bx(ei)-x0
        by = rt + (ROW_H-BAR_H)//2
        sp = add_rect(sg, x0+12000, by, w-24000, BAR_H, fill, line_color=tcol, line_w=0.75)
        tf = sp.text_frame; tf.word_wrap=True
        tf.margin_left=Emu(27000); tf.margin_right=Emu(27000); tf.margin_top=0; tf.margin_bottom=0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=txt; set_font(r, 7, tcol, bold=True)
# legend swatches
lx = 571500; ly = CHART_B + 200000
for fill_pair, lab in [(DONE,"완료 (현재까지)"),(S1c,"Sprint 1 사전개발 (6/5~6/12)"),(S2c,"Sprint 2 본선현장 (6/13~14)")]:
    add_rect(sg, lx, ly+15000, 200000, 150000, fill_pair[0], line_color=fill_pair[1], line_w=0.75)
    tb = textbox(sg, lx+250000, ly, 2700000, 200000, lab, 8.5, BODY)
    lx += 3050000
textbox(sg, lx, ly, 3000000, 200000, "◆ 마일스톤 (예선결과·본선·발표)", 8.5, NAVY, bold=True)

# ============== SLIDE 3: 상세 개발 계획 (30분 단위) ==============
s2 = add_slide()
header(s2,
       "본선 무박 2일을 30분 단위로 설계한 실행 계획",
       "Sprint 1 사전 개발(일자별) + Sprint 2 본선 현장(30분 블록) — 시간 손실 없는 실행 로드맵")

# Sprint 1 strip (full width table 2 rows x 5)
s1cols = ["6/5 킥오프", "6/6~6/7 데이터", "6/8~6/9 AI 코어", "6/10~6/11 카카오·게임", "6/12 통합"]
s1vals = ["예선결과 확인·역할 재확정·repo 전략",
          "좌표/행사 연동·bge-m3 전체 재임베딩·alias 100+",
          "6브랜치 라우터·Intent·Soft Filter·프롬프트 튜닝",
          "Skill 응답포맷·멀티턴·퀴즈/카드/도감 WebView",
          "E2E 통합 테스트·버그수정·배포 확정"]
textbox(s2, 571500, 1330000, 5000000, 220000, "Sprint 1 — 사전 개발 (6/5~6/12)", 11, NAVY, bold=True)
t1 = make_table(s2, 2, 5, 571500, 1600000, 11049000, 1150000, [2209800]*5)
for j, c in enumerate(s1cols):
    style_cell(t1.cell(0, j), c, 9, WHITE, NAVY, bold=True, align=PP_ALIGN.CENTER)
for j, v in enumerate(s1vals):
    style_cell(t1.cell(1, j), v, 8, G1[1], G1[0], align=PP_ALIGN.CENTER)

# Sprint 2 — Day1 (left) / Day2 (right), 30min blocks
day1 = [
 ("11:30~12:00", "팀 킥오프 · Sprint 2 목표/역할 재확인"),
 ("13:00~13:30", "배포 상태 점검 · 현장 서버 확인"),
 ("13:30~14:30", "말풍선 글자수 제어 · 온보딩 웰컴 메시지"),
 ("14:30~15:30", "멀티턴 컨텍스트 유지 · RAG 품질 점검"),
 ("15:30~16:00", "중간 통합 테스트"),
 ("16:00~17:00", "초기화 커맨드 · 퀴즈 세션 관리"),
 ("17:00~18:00", "모르는 유산 fallback · 발표자료 착수"),
 ("19:00~20:30", "이미지 카드(Carousel) · 피드백 수집 · 발표자료"),
 ("21:00~21:30", "Day 1 회고 · Day 2 우선순위 재조정"),
 ("21:30~24:00", "비교질문(P1) · 엣지케이스 · 응답 2차 튜닝"),
]
day2 = [
 ("00:00~02:00", "P1 잔여 기능 · 전체 최종 통합 테스트"),
 ("03:00~05:00", "서버 안정화·부하 테스트 · 발표자료 마무리"),
 ("06:00~07:00", "조식 · 최종 서버 상태 확인"),
 ("07:00~08:30", "최종 버그 수정 · 발표자료 완성"),
 ("08:30~09:30", "데모 시나리오 리허설(경복궁→해설→퀴즈→추천)"),
 ("09:30~10:30", "발표 리허설 1회(시간 측정)"),
 ("10:30~11:30", "피드백 반영 · 데모 계정/서버 최종 점검"),
 ("13:00~", "발표 + 심사·질의응답 → 결과 발표"),
]
textbox(s2, 571500, 2950000, 5000000, 220000, "Sprint 2 — 본선 6/13 (Day 1)", 11, NAVY, bold=True)
textbox(s2, 6150000, 2950000, 5000000, 220000, "Sprint 2 — 본선 6/14 (Day 2)", 11, NAVY, bold=True)

def time_table(slide, data, l):
    t = make_table(slide, len(data)+1, 2, l, 3220000, 5470000, 3000000, [1250000, 4220000])
    style_cell(t.cell(0,0), "시간", 8.5, WHITE, NAVY, bold=True, align=PP_ALIGN.CENTER)
    style_cell(t.cell(0,1), "작업", 8.5, WHITE, NAVY, bold=True)
    for i,(tm,work) in enumerate(data, start=1):
        style_cell(t.cell(i,0), tm, 7.5, G2[1], G2[0], bold=True, align=PP_ALIGN.CENTER)
        style_cell(t.cell(i,1), work, 7.5, BODY, WHITE)
    return t
time_table(s2, day1, 571500)
time_table(s2, day2, 6150000)

textbox(s2, 571500, 6350000, 11049000, 200000,
        "※ 공식 현장 일정은 별도 고지 전이나, 30분 단위 자체 계획으로 무박 2일의 실행력을 확보 (P0 우선, 시간 부족 시 compare→이미지 순으로 컷)",
        8.5, GRAY)

prs.save(DST)
print("saved:", DST, "| total slides:", len(prs.slides._sldIdLst))
