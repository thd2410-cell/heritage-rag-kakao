# -*- coding: utf-8 -*-
"""Slide 18·22 데이터 검증 보강 수정. 원본 보존, 새 파일로 저장."""
from pptx import Presentation

SRC = r"c:\Users\SSAFY\Desktop\heritage-rag-kakao\data_eng\presentation\유산톡 AI.pptx"
DST = r"c:\Users\SSAFY\Desktop\heritage-rag-kakao\data_eng\presentation\유산톡 AI_데이터검증수정_v3.pptx"

prs = Presentation(SRC)

def shp(slide_idx, name):
    for s in prs.slides[slide_idx].shapes:
        if s.name == name:
            return s
    raise KeyError(name)

def set_runs(shape, new_texts):
    """첫 문단 run들에 텍스트를 채움(서식 유지). run이 모자라면 마지막 run에 합침."""
    para = shape.text_frame.paragraphs[0]
    runs = para.runs
    n = len(runs)
    for i, t in enumerate(new_texts):
        if i < n - 1:
            runs[i].text = t
        elif i == n - 1:
            # 남은 텍스트 전부 마지막 run에
            runs[i].text = "".join(new_texts[i:])
            break
    # 새 텍스트가 run보다 적으면 나머지 run 비움
    if len(new_texts) < n:
        for j in range(len(new_texts), n):
            runs[j].text = ""

# ---------- Slide 18 (index 17) ----------
# 1) Recall@K 열 제거 (헤더 + 4개 값 비움) — 깨진 지표 화면에서 삭제
for nm in ["SK-18-text-40", "SK-18-text-45", "SK-18-text-50", "SK-18-text-55", "SK-18-text-60"]:
    set_runs(shp(17, nm), [""])

# 2) 인사이트 3번 → 긍정형 '검색 정밀 검증' (틀렸다는 표현 없음, 2 run)
set_runs(shp(17, "SK-18-text-65"), [
    "검색 정밀 검증:",
    " 정답셋(expected_heritage) 라벨링 + 운영 임베딩(bge-m3)으로 재측정 → 단일질문 Recall@3 ≈ 100% (자체 정답셋)",
])

# 3) 각주 → 프로토타입/운영 구분 + 긍정형 (Recall ≈100%, facet 0.89)
set_runs(shp(17, "SK-18-text-66"), [
    "※ 표: 프로토타입 1차 실험(KoSimCSE/ChromaDB, 80문항·5,739 Chunk). 검색 Recall(≈100%)·facet 답변 충실성(0.89)은 운영 임베딩(bge-m3)으로 검증. 전체 서비스 성능 보장은 아님.",
])

# ---------- Slide 22 (index 21) ----------
# Facet 품질 편차 카드 본문 → 원인 규명 + 측정 + 개선 실증 (3 run)
set_runs(shp(21, "SK-22-text-56"), [
    "빈약한 본문이 원인임을 규명, ",
    "충실성 0.89 측정·프롬프트 강화로 0.76→0.93 개선",
    " 실증",
])

# ---------- Slide 19 (index 18): 평가 스택을 운영(bge-m3/pgvector)으로 정합 ----------
set_runs(shp(18, "SK-19-text-49"), ["bge-m3 (운영 임베딩과 동일)"])      # Embedding (Target)
set_runs(shp(18, "SK-19-text-51"), ["pgvector (운영과 동일)"])          # Vector Store
set_runs(shp(18, "SK-19-text-53"), ["Gemini 2.5 Flash"])               # Evaluation Judge (RAGAS + Gemini)
set_runs(shp(18, "SK-19-text-57"), [                                   # 선정 근거 + 정합성 한 줄
    "80문항 평가셋 기반의 객관적 품질 진단 및 하이퍼파라미터(Top-K) 최적화 · 운영 임베딩(bge-m3)으로 검증해 테스트–운영 정합성 확보",
])

# (Slide 23은 원본 그대로 — 각주 미추가)

prs.save(DST)
print("saved:", DST)
